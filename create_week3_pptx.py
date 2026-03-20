#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

OUTPUT_PATH = "/mnt/c/Users/miran/Automation Workz Dropbox/Miranda's Workspace/Miranda Boyd/IoT Professional Class/presentations/Week3_Presentation.pptx"


def add_slide(pres, title, content_items=None, subtitle=None, is_quote=False):
    slide = pres.slides.add_slide(pres.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 102, 0)

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        for para in subtitle_shape.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.italic = True

    if content_items:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 49, 68)
            p.level = 0

    return slide


def create_week3_presentation():
    pres = Presentation()
    pres.slide_width = Inches(10)
    pres.slide_height = Inches(7.5)

    slides_content = [
        {
            "title": "Week 3: Financial Gap Analysis & Business Planning",
            "subtitle": "Automation Workz | IoT Professional Class",
        },
        {
            "title": "From Vision to Reality",
            "subtitle": '"A gap analysis without financial data is just a wish list."',
        },
        {
            "title": "Strategic Takeaway",
            "content_items": [
                '"A gap analysis without financial data is just a wish list.',
                "",
                "When you connect your VISION to your NUMBERS,",
                'you create a roadmap that can actually be executed."',
            ],
        },
        {
            "title": "Learning Objectives",
            "content_items": [
                "1. Upload and analyze financial statements",
                "2. Connect financial data to Week 2 Vision",
                "3. Identify specific gaps between current state and goals",
                "4. Create a Gap-to-Goal Action Plan",
                "5. Use autoworkz.org for financial analysis",
            ],
        },
        {
            "title": "Your 4 Data Points",
            "content_items": [
                "OCEAN Profile (Week 1) - Who you are naturally",
                "Vision Board (Week 2) - Where you want to go",
                "Tech Readiness (Week 2) - How ready you are for change",
                "Financial Reality (Week 3) - Where you stand today",
            ],
        },
        {
            "title": "Three Key Financial Statements",
            "content_items": [
                "Income Statement (P&L) - Revenue, expenses, profit/loss over time",
                "Balance Sheet - Assets, liabilities, equity at a point in time",
                "Cash Flow Statement - Money moving in and out over time",
            ],
        },
        {
            "title": "Types of Gaps",
            "content_items": [
                "Revenue Gap - Current vs. target revenue ($100K vs. $250K = $150K gap)",
                "Customer Gap - Current vs. ideal customer profile",
                "Resource Gap - Available vs. needed resources",
                "Time Gap - Current vs. desired timeline",
            ],
        },
        {
            "title": "In-Class Activities",
            "content_items": [
                "Step 1: Prepare Financial Documents (~15 min)",
                "Step 2: Upload to autoworkz.org (~25 min)",
                "Step 3: Map Reality to Vision (~20 min)",
                "Step 4: Create Gap-to-Goal Action Plan (~30 min)",
            ],
        },
        {
            "title": "Step 1: Prepare Financial Documents",
            "content_items": [
                "Gather your three key financial statements:",
                "   - Income Statement (P&L)",
                "   - Balance Sheet",
                "   - Cash Flow Statement",
                "Ensure documents are current (within last 12 months)",
                "Have digital copies ready for upload",
            ],
        },
        {
            "title": "Step 2: Upload to autoworkz.org",
            "content_items": [
                "Log into autoworkz.org",
                "Upload your financial documents",
                "Review the AI-generated analysis",
                "Note key metrics and insights",
                "",
                "Your data is private and secure",
            ],
        },
        {
            "title": "Step 3: Map Reality to Vision",
            "content_items": [
                "Compare financials with Week 2 Vision Board",
                "",
                "Revenue Gap: Current revenue vs. target?",
                "Customer Gap: Current vs. ideal profile?",
                "Resource Gap: Do you have funds to hire/expand?",
                "Time Gap: How long to reach goals?",
            ],
        },
        {
            "title": "Step 4: Gap-to-Goal Action Plan",
            "content_items": [
                "Current State: Where you are now",
                "Goal State: Where you want to be",
                "Gap: Measured in $, time, or resources",
                "Actions: 3-5 specific steps per gap",
                "Timeline: When will each action complete?",
            ],
        },
        {
            "title": "Key Takeaway",
            "subtitle": '"Connect your vision to your numbers to create an executable roadmap."',
            "content_items": [
                "",
                "Automation Workz | IoT Professional Class",
                "440 Burroughs St. Ste. 173, Detroit, MI 48202",
                "313-483-2126 | ida@autoworkz.org",
            ],
        },
    ]

    for slide_data in slides_content:
        add_slide(
            pres,
            slide_data.get("title", ""),
            slide_data.get("content_items"),
            slide_data.get("subtitle"),
        )

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    pres.save(OUTPUT_PATH)
    print(f"Created: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    create_week3_presentation()
