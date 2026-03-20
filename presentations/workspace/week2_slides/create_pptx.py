#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

COLORS = {
    "dark_purple": RGBColor(0x0F, 0x05, 0x68),
    "orange": RGBColor(0xFF, 0x66, 0x00),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark_blue": RGBColor(0x00, 0x00, 0x66),
    "dark_gray": RGBColor(0x33, 0x33, 0x33),
    "light_gray": RGBColor(0xF5, 0xF5, 0xF5),
    "cream": RGBColor(0xFF, 0xF5, 0xEE),
}


def add_header_bar(slide, title_text):
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.85)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["dark_purple"]
    header.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.13), Inches(0.1), Inches(0.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["orange"]
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]


def add_time_badge(slide, text, left, top):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.2), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS["orange"]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]

# Slide 1: Title
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.solid()
bg.fill.fore_color.rgb = COLORS["dark_purple"]
bg.line.fill.background()

week_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(3), Inches(1.5), Inches(4), Inches(0.5)
)
week_bar.fill.solid()
week_bar.fill.fore_color.rgb = COLORS["orange"]
week_bar.line.fill.background()
tf = week_bar.text_frame
p = tf.paragraphs[0]
p.text = "WEEK 2"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Business Vision & Tech Readiness"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(3.3), Inches(9), Inches(0.5)
)
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = '"Transform Self-Awareness into Strategic Action"'
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = COLORS["orange"]
p.alignment = PP_ALIGN.CENTER

# Slide 2: Week 1 Homework Review
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Week 1 Homework Review: What You Learned")

subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.0), Inches(9), Inches(0.35)
)
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "You used the Executive Personality Prompt with your OCEAN scores:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(4.3), Inches(2.5))
tf = left_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "What AI Revealed About You:"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS["dark_blue"]
insights = [
    "Executive strengths & blind spots",
    "Leadership tendencies",
    "Strategic decision-making style",
    "Marketing & branding approach",
    "Problem-solving & innovation style",
]
for item in insights:
    p = tf.add_paragraph()
    p.text = f"  - {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["dark_gray"]
    p.space_before = Pt(4)

right_box = slide.shapes.add_textbox(Inches(5), Inches(1.45), Inches(4.5), Inches(2.5))
tf = right_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "5 OCEAN Traits Mapped To:"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS["dark_blue"]
mappings = [
    "Leadership/Team Style",
    "Strategic Thinking",
    "Marketing/Communication Style",
]
for item in mappings:
    p = tf.add_paragraph()
    p.text = f"  - {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["dark_gray"]
    p.space_before = Pt(4)

question_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.3), Inches(9), Inches(0.6)
)
tf = question_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Discussion: How did your OCEAN traits reveal blind spots in your leadership or marketing?"
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = COLORS["orange"]

# Slide 3: Activity 1
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Week 2: Activity 1")

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Business Vision Board"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(2.5))
tf = content_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Using SenseiiWyze (senseiiwyze.com), create a comprehensive vision board that captures your business aspirations and personal strengths."
p.font.size = Pt(14)
p.font.color.rgb = COLORS["dark_gray"]
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "Your vision board will help you see the connection between who you are (OCEAN profile), where you want to go (Vision), and how ready you are for change (Tech Readiness)."
p.font.size = Pt(14)
p.font.color.rgb = COLORS["dark_gray"]

add_time_badge(slide, "25 minutes", 0.5, 4.2)

# Slide 4: 8 Key Elements
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "8 Key Elements of Your Vision Board")

elements = [
    ("1. Personal Identity", "Who you are as a leader"),
    ("2. Core Values", "What principles guide you"),
    ("3. Vision Alignment", "Your 3-5 year business goals"),
    ("4. Customer Persona", "Who you serve best"),
    ("5. Brand Identity", "How you want to be known"),
    ("6. Dream Team", "Who you need beside you"),
    ("7. Leadership Skills", "Strengths to develop"),
    ("8. Metrics & Growth", "How you measure success"),
]
positions = [
    (0.4, 1.1),
    (2.8, 1.1),
    (5.2, 1.1),
    (7.6, 1.1),
    (0.4, 2.0),
    (2.8, 2.0),
    (5.2, 2.0),
    (7.6, 2.0),
]

for (title, desc), (x, y) in zip(elements, positions):
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.2), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["light_gray"]
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark_blue"]
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["dark_gray"]

# Slide 5: Activity 2
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Week 2: Activity 2")

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Technical Skills Assessment"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.4))
tf = content_box.text_frame
p = tf.paragraphs[0]
p.text = "Complete 5 levels of tech skill games at Senseii Games (senseiigames.com)"
p.font.size = Pt(14)
p.font.color.rgb = COLORS["dark_gray"]

levels = [
    "Level 1: Basic Digital Navigation",
    "Level 2: Data Entry & Management",
    "Level 3: Communication Tools",
    "Level 4: Business Applications",
    "Level 5: Advanced Problem Solving",
]
list_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(2))
tf = list_box.text_frame
for i, level in enumerate(levels):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = level
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["dark_gray"]
    p.space_before = Pt(6)

add_time_badge(slide, "20 minutes", 0.5, 4.2)

# Slide 6: Why Tech Readiness Matters
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Why Tech Readiness Matters")

stats = [
    ("10X", "Training ROI with tech-enabled coaching"),
    ("58%", "Higher employee retention rate"),
    ("2.5X", "Faster promotion rate for tech-ready leaders"),
]
for i, (stat, desc) in enumerate(stats):
    x = 0.5 + i * 3.1
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(2.8), Inches(1.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["light_gray"]
    box.line.fill.background()
    stat_box = slide.shapes.add_textbox(
        Inches(x), Inches(1.4), Inches(2.8), Inches(0.6)
    )
    tf = stat_box.text_frame
    p = tf.paragraphs[0]
    p.text = stat
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    desc_box = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(2.0), Inches(2.6), Inches(0.7)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["dark_gray"]
    p.alignment = PP_ALIGN.CENTER

quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Tech readiness is not just about skills - it is about adaptability in a changing business landscape."
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = COLORS["dark_blue"]
p.alignment = PP_ALIGN.CENTER

# Slide 7: 3 Data Points
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Your Complete Business Picture: 3 Data Points")

subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.1), Inches(9), Inches(0.4)
)
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "When combined, these reveal your complete business profile:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

data_points = [
    ("OCEAN Profile", "Who you are naturally"),
    ("Vision Board", "Where you want to go"),
    ("Tech Readiness", "How ready you are for change"),
]
for i, (title, desc) in enumerate(data_points):
    x = 0.3 + i * 2.4
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.7), Inches(2.2), Inches(1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["dark_purple"]
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

equals_box = slide.shapes.add_textbox(
    Inches(7.5), Inches(1.9), Inches(0.5), Inches(0.5)
)
tf = equals_box.text_frame
p = tf.paragraphs[0]
p.text = "="
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]
p.alignment = PP_ALIGN.CENTER

result_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.7), Inches(1.7), Inches(1)
)
result_box.fill.solid()
result_box.fill.fore_color.rgb = COLORS["orange"]
result_box.line.fill.background()
tf = result_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Complete Picture"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Your strategic business identity"
p.font.size = Pt(9)
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

# Slide 8: Finding & Fixing Gaps
slide = prs.slides.add_slide(blank_layout)
add_header_bar(slide, "Finding & Fixing the Gaps")

subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.0), Inches(9), Inches(0.4)
)
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Common gaps between vision and reality:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS["orange"]

gaps = [
    ("Skills Gap", "Vision requires skills you have not developed yet"),
    ("Personality Gap", "Your OCEAN traits may conflict with your goals"),
    ("Tech Gap", "Your tech readiness is below what your vision demands"),
    ("Resource Gap", "Missing team, funding, or tools needed"),
]
for i, (title, desc) in enumerate(gaps):
    x = 0.4 + i * 2.4
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(2.2), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["cream"]
    box.line.color.rgb = COLORS["orange"]
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["dark_gray"]
    p.alignment = PP_ALIGN.CENTER

# Slide 9: Key Takeaway
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.solid()
bg.fill.fore_color.rgb = COLORS["dark_purple"]
bg.line.fill.background()

week_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(3), Inches(1), Inches(4), Inches(0.4)
)
week_bar.fill.solid()
week_bar.fill.fore_color.rgb = COLORS["orange"]
week_bar.line.fill.background()
tf = week_bar.text_frame
p = tf.paragraphs[0]
p.text = "Week 2: Key Takeaway"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(0.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Strategic Takeaway"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

quote_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(2)
)
quote_bg.fill.solid()
quote_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x7A)
quote_bg.line.fill.background()

quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(7.6), Inches(1.6))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"The gap between where you are and where you want to be is not a problem - it is your roadmap. Your OCEAN profile, Vision Board, and Tech Readiness score together reveal exactly what bridges you need to build."'
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = COLORS["white"]
p.alignment = PP_ALIGN.CENTER

output_path = "/mnt/c/Users/miran/Automation Workz Dropbox/Miranda's Workspace/Miranda Boyd/IoT Professional Class/presentations/Week2_InClassActivities_Aligned.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
